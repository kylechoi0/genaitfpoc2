import requests
import streamlit as st
import json
import traceback
import PyPDF2
import io
import docx2txt
import pandas as pd
from pptx import Presentation
import time

# 별도의 API 키 설정 (전처리 워크플로우용)
PREPROCESS_API_KEY = st.secrets["PREPROCESS_API_KEY"]
PREPROCESS_API_URL = 'https://mir-api.52g.ai/v1/workflows/run'
KNOWLEDGE_API_KEY = st.secrets["KNOWLEDGE_API_KEY"]

def extract_text_from_file(uploaded_file):
    file_extension = uploaded_file.name.lower().split('.')[-1]
    
    try:
        if file_extension == 'pdf':
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(uploaded_file.getvalue()))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
                
        elif file_extension in ['doc', 'docx']:
            text = docx2txt.process(io.BytesIO(uploaded_file.getvalue()))
            
        elif file_extension == 'txt':
            text = uploaded_file.getvalue().decode('utf-8')
            
        elif file_extension == 'md':
            text = uploaded_file.getvalue().decode('utf-8')
            
        elif file_extension in ['ppt', 'pptx']:
            prs = Presentation(io.BytesIO(uploaded_file.getvalue()))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif file_extension in ['xls', 'xlsx', 'csv']:
            if file_extension == 'csv':
                df = pd.read_csv(uploaded_file)
            else:
                df = pd.read_excel(uploaded_file)
            text = df.to_string()
            
        elif file_extension == 'hwp':
            st.error("HWP 파일 형식은 현재 지원되지 않습니다. PDF로 변환 후 시도해주세요.")
            return None
            
        else:
            st.error(f"지원하지 않는 파일 형식입니다: {file_extension}")
            return None
            
        return text.strip()
        
    except Exception as e:
        st.error(f"파일 처리 중 오류 발생: {str(e)}")
        return None

def preprocess_files(uploaded_files, dataset_id):
    headers = {
        'Authorization': f'Bearer {PREPROCESS_API_KEY}',
    }

    try:
        # 파일 크기 검사
        file = uploaded_files[0]
        if file.size > 200 * 1024 * 1024:  # 200MB
            st.error("파일 크기가 200MB를 초과합니다.")
            return None
            
        # 텍스트 추출
        extracted_text = extract_text_from_file(file)
        if not extracted_text:
            return None
            
        # 워크플로우 실행 요청
        workflow_payload = {
            'response_mode': 'blocking',
            'user': 'user-123',
            'inputs': {
                'text': extracted_text
            },
            'workflow_id': '6a157fa1-8f3d-4bde-8d8c-78df231a724c'
        }
        
        # 진행 상태 표시
        progress_text = "파일 처리 중..."
        status_container = st.empty()
        status_container.text(progress_text)
        
        # 단일 요청으로 처리 (10분 타임아웃)
        workflow_response = requests.post(
            PREPROCESS_API_URL,
            headers=headers,
            json=workflow_payload,
            timeout=600  # 10분 타임아웃
        )
        
        if workflow_response.status_code == 200:
            status_container.text("처리 완료!")
            result = workflow_response.json()
            file_url = result.get('data', {}).get('outputs', {}).get('result')
            
            # 지식 데이터셋에 문서 추가
            knowledge_headers = {
                'Authorization': f'Bearer {KNOWLEDGE_API_KEY}',
                'Content-Type': 'application/json'
            }
            
            knowledge_payload = {
                'name': file.name,
                'text': extracted_text,
                'indexing_technique': 'high_quality',
                'process_rule': {
                    'mode': 'custom',
                    'rules': {
                        'pre_processing_rules': [
                            {'id': 'remove_extra_spaces', 'enabled': True},
                            {'id': 'remove_urls_emails', 'enabled': True}
                        ],
                        'segmentation': {
                            'separator': '####',
                            'max_tokens': 1000
                        }
                    }
                }
            }
            
            knowledge_response = requests.post(
                f'https://mir-api.52g.ai/v1/datasets/{dataset_id}/document/create_by_text',
                headers=knowledge_headers,
                json=knowledge_payload
            )
            
            # 다운로드 링크 추가 (form 밖으로 이동)
            if file_url:
                st.markdown("### 전처리된 파일 다운로드")
                processed_filename = file.name.rsplit('.', 1)[0] + '_processed.txt'
                st.markdown(f'<a href="{file_url}" download="{processed_filename}" target="_blank">📥 전처리된 파일 다운로드</a>', unsafe_allow_html=True)
            
            if knowledge_response.status_code == 200:
                st.success("지식 데이터셋에 문서가 추가되었습니다!")
                return knowledge_response.json()
            else:
                st.error(f"지식 데이터셋 추가 실패: {knowledge_response.text}")
                return None
                
        else:
            status_container.error("처리 실패")
            st.error(f"상태 코드: {workflow_response.status_code}")
            st.error(f"응답 내용: {workflow_response.text}")
            return None

    except requests.exceptions.Timeout:
        st.error("처리 시간이 초과되었습니다 (10분). 잠시 후 다시 시도해주세요.")
        return None
    except Exception as e:
        st.error(f"⚠️ 처리 중 오류 발생: {str(e)}")
        st.error(traceback.format_exc())
        return None

def upload_to_knowledge_directly(file, dataset_id):
    try:
        headers = {
            'Authorization': f'Bearer {KNOWLEDGE_API_KEY}'
        }
        
        # multipart/form-data 형식으로 데이터 준비
        files = {
            'file': (file.name, file, 'application/octet-stream')
        }
        
        # 자동 처리 설정
        data = {
            'data': json.dumps({
                'indexing_technique': 'high_quality',
                'process_rule': {
                    'mode': 'automatic'  # 'auto'가 아닌 'automatic'으로 수정
                }
            })
        }
        
        # 파일 업로드 요청
        response = requests.post(
            f'https://mir-api.52g.ai/v1/datasets/{dataset_id}/document/create_by_file',
            headers=headers,
            files=files,
            data=data,
            timeout=30
        )
        
        if response.status_code == 200:
            result = response.json()
            return result
        else:
            st.error(f"파일 업로드 실패 (상태 코드: {response.status_code})")
            st.error(f"오류 내용: {response.text}")
            return None
            
    except requests.exceptions.Timeout:
        st.error("업로드 시간이 초과되었습니다. 잠시 후 다시 시도해주세요.")
        return None
    except Exception as e:
        st.error(f"파일 처리 중 오류 발생: {str(e)}")
        return None
